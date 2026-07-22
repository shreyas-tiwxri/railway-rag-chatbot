"""Loaders for non-PDF document types: Word, PowerPoint, Excel, and web pages.
Each returns the same page-like structure as pdf_loader.extract_pages_text:
[{"page_number": int, "text": str, "ocr_used": False}, ...]
so the rest of the ingestion pipeline (chunking, embedding) doesn't need to
know or care what the original file format was.
"""
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import requests
from bs4 import BeautifulSoup


def extract_docx_pages(path: str) -> list[dict]:
    doc = DocxDocument(path)
    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    # Word docs don't have a reliable page concept via python-docx, so the
    # whole document is treated as one "page" — the chunker still splits it
    # into multiple embeddable chunks, just without a meaningful page number
    # (citations for Word docs will say "page 1" regardless of true location).
    return [{"page_number": 1, "text": full_text, "ocr_used": False}]


def extract_pptx_pages(path: str) -> list[dict]:
    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        pages.append({"page_number": i + 1, "text": "\n".join(texts), "ocr_used": False})
    return pages


def extract_xlsx_pages(path: str) -> list[dict]:
    wb = openpyxl.load_workbook(path, data_only=True)
    pages = []
    for i, sheet_name in enumerate(wb.sheetnames):
        sheet = wb[sheet_name]
        rows_text = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows_text.append(" | ".join(cells))
        text = f"Sheet: {sheet_name}\n" + "\n".join(rows_text)
        pages.append({"page_number": i + 1, "text": text, "ocr_used": False})
    return pages


def extract_url_page(url: str) -> list[dict]:
    resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=20)
    resp.raise_for_status()
    soup = BeautifulSoup(resp.text, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = " ".join(soup.get_text(separator=" ").split())
    return [{"page_number": 1, "text": text, "ocr_used": False}]
